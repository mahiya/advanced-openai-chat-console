import os
from flask import Flask, Response, request # Flask
from utilities.multithread_openai import get_completion, get_completion_from_long_text # OpenAI Client
from utilities.tmpfile import tmpfile # 一時ファイル処理
from pdfminer.high_level import extract_text # PDFファイル処理
import docx # Microsoft Word ファイル処理
from pptx import Presentation # Microsoft Power Point ファイル処理

app = Flask(__name__)

@app.route("/", defaults={"path": "index.html"})
@app.route("/<path:path>")
def static_file(path):
    return app.send_static_file(path)

@app.route("/conversation", methods=["POST"])
def conversation():
    message = request.json['message']
    history = request.json['history']
    completion = get_completion(message, history)
    return completion, 200

@app.route("/upload", methods=["POST"])
def get_completion_from_file():
    file_name = request.form["fileName"]
    file_path, file_ext = os.path.splitext(file_name)
    if file_ext == '.txt': return get_completion_from_txt()
    elif file_ext == '.pdf': return get_completion_from_pdf()
    elif file_ext == '.docx': return get_completion_from_docx()
    elif file_ext == '.pptx': return get_completion_from_pptx()
    else: return "Uploaded file type is not supported.", 400

# テキストファイル(.txt)の処理
def get_completion_from_txt():
    file_bytes = request.files['file'].stream.read()
    text = file_bytes.decode('utf-8')
    return get_completion_from_long_text(request.form["message"], text), 200

# PDFファイル(.pdf)の処理
def get_completion_from_pdf():
    file_bytes = request.files['file'].stream.read()
    with tmpfile(file_bytes) as f:
        text = extract_text(f.tmp_file_name).replace("\n", "。")
    return get_completion_from_long_text(request.form["message"], text), 200

# Microsoft Word ファイル(.docx)の処理
def get_completion_from_docx():
    file_bytes = request.files['file'].stream.read()
    with tmpfile(file_bytes) as f:
        doc = docx.Document(f.tmp_file_name)
        text = " ".join([p.text for p in doc.paragraphs])
    return get_completion_from_long_text(request.form["message"], text), 200

# Microsoft Power Point ファイル(.pptx)の処理
def get_completion_from_pptx():
    file_bytes = request.files['file'].stream.read()
    with tmpfile(file_bytes) as f:
        pr = Presentation(f.tmp_file_name)
        texts = []
        for slide in pr.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                text = shape.text_frame.text.strip()
                if len(text) > 0: texts.append(text)
        text = " ".join(texts)
    return get_completion_from_long_text(request.form["message"], text), 200

if __name__ == "__main__":
    app.run()